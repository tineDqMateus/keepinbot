"""
create_samples.py
Génération des fichiers de test pour le Module 2
Formats : Word (.docx), PowerPoint (.pptx), mail (.eml)
"""

import os
from docx import Document
from pptx import Presentation
from pptx.util import Inches, Pt

os.makedirs("data/samples", exist_ok=True)

# ── Fichier Word — compte-rendu de réunion ────────────────────────────────────
doc = Document()
doc.add_heading("Compte-rendu réunion projet Alpha", 0)
doc.add_paragraph("Date : 15 mars 2024")
doc.add_paragraph("Participants : Marie Dupont (chef de projet), Jean Martin (dev), Sophie Leblanc (RH)")
doc.add_heading("Décisions prises", level=1)
doc.add_paragraph("Le projet Alpha démarrera le 1er avril 2024 pour une durée de 6 mois.")
doc.add_paragraph("Budget alloué : 45 000 euros. Validation requise par la direction avant le 20 mars.")
doc.add_paragraph("3 sessions de formation prévues en avril, mai et juin pour les équipes concernées.")
doc.add_heading("Points ouverts", level=1)
doc.add_paragraph("Recrutement d'un développeur supplémentaire en cours — réponse attendue avant fin mars.")
doc.add_paragraph("Choix de l'outil de gestion de projet non finalisé : arbitrage entre Jira et Notion.")
doc.add_heading("Prochaine réunion", level=1)
doc.add_paragraph("22 mars 2024 à 14h — salle de réunion B.")
doc.save("data/samples/cr_reunion_alpha.docx")
print("Créé : data/samples/cr_reunion_alpha.docx")

# ── Fichier PowerPoint — présentation projet ──────────────────────────────────
prs = Presentation()

slide1 = prs.slides.add_slide(prs.slide_layouts[0])
slide1.shapes.title.text = "Projet Alpha — Lancement"
slide1.placeholders[1].text = "Présentation équipe — Mars 2024"

slide2 = prs.slides.add_slide(prs.slide_layouts[1])
slide2.shapes.title.text = "Contexte et objectifs"
slide2.placeholders[1].text = "Moderniser le système de gestion des commandes.\nRéduire les délais de traitement de 30%.\nDéploiement prévu : octobre 2024."

slide3 = prs.slides.add_slide(prs.slide_layouts[1])
slide3.shapes.title.text = "Planning"
slide3.placeholders[1].text = "Avril : cadrage technique et recrutement.\nMai-Juin : développement phase 1.\nJuillet : tests et recette.\nOctobre : mise en production."

slide4 = prs.slides.add_slide(prs.slide_layouts[1])
slide4.shapes.title.text = "Budget et ressources"
slide4.placeholders[1].text = "Budget total : 45 000 euros.\nÉquipe : 3 développeurs, 1 chef de projet, 1 UX designer.\nOutils : Jira pour le suivi, GitHub pour le code."

prs.save("data/samples/presentation_alpha.pptx")
print("Créé : data/samples/presentation_alpha.pptx")

# ── Fichier mail — échange sur le projet ─────────────────────────────────────
mail_content = """From: marie.dupont@artisanplus.fr
To: jean.martin@artisanplus.fr
Subject: RE: Projet Alpha — point recrutement
Date: 18 mars 2024

Jean,

Suite à notre réunion du 15 mars, voici les points importants à retenir :

Le recrutement du développeur supplémentaire est confirmé. Le candidat retenu
est Paul Durand, disponible à partir du 2 avril. Son profil Python/Django
correspond exactement à nos besoins sur le projet Alpha.

Concernant l'outil de gestion de projet, la décision a été prise par la direction :
nous utiliserons Jira. Les licences seront commandées cette semaine.

Le budget de 45 000 euros est validé. Tu peux démarrer les achats de matériel
dès la semaine prochaine. Pense à conserver toutes les factures pour la
comptabilité.

Prochaine étape : réunion de lancement officiel le 1er avril à 9h, salle A.
Toute l'équipe doit être présente.

Cordialement,
Marie Dupont
Chef de projet — Artisan Plus SARL
"""

with open("data/samples/mail_recrutement.eml", "w", encoding="utf-8") as f:
    f.write(mail_content)
print("Créé : data/samples/mail_recrutement.eml")

print("\nCorpus de test généré — 3 fichiers dans data/samples/")