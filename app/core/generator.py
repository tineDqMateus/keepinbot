"""
generator.py
Module 2 — Générateur de documentation Keepinbot
=================================================
Deux responsabilités :
1. Parsers multi-format : extraire le texte de fichiers hétérogènes
   (Word, PowerPoint, mail, PDF, texte brut)
2. Génération : à partir des fragments extraits, produire un document
   structuré via le LLM

Principe de validation :
Aucun document généré n'est automatiquement indexé dans la base RAG.
L'utilisateur doit valider explicitement avant indexation.
"""

import os
import email
from docx import Document as DocxDocument
from pptx import Presentation
from app.core.config import OLLAMA_BASE_URL, OLLAMA_MODEL, MISTRAL_API_KEY, MISTRAL_MODEL, LLM_MODE
import requests


# ── Parsers multi-format ──────────────────────────────────────────────────────

def parse_txt(filepath: str) -> str:
    """
    Lit un fichier texte brut (.txt).

    Paramètre :
    - filepath : chemin complet vers le fichier

    Retourne :
    - contenu texte du fichier (str)
    """
    with open(filepath, "r", encoding="utf-8") as f:
        return f.read()


def parse_docx(filepath: str) -> str:
    """
    Extrait le texte d'un fichier Word (.docx).
    Lit les paragraphes dans l'ordre — titres et corps de texte.
    Les tableaux et images sont ignorés.

    Paramètre :
    - filepath : chemin complet vers le fichier .docx

    Retourne :
    - texte extrait, paragraphes séparés par des sauts de ligne (str)
    """
    doc = DocxDocument(filepath)
    paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]
    return "\n".join(paragraphs)


def parse_pptx(filepath: str) -> str:
    """
    Extrait le texte d'un fichier PowerPoint (.pptx).
    Lit le contenu de chaque diapositive dans l'ordre.
    Préfixe chaque diapositive par son numéro pour conserver
    la structure de la présentation.

    Paramètre :
    - filepath : chemin complet vers le fichier .pptx

    Retourne :
    - texte extrait diapositive par diapositive (str)
    """
    prs = Presentation(filepath)
    slides_text = []
    for i, slide in enumerate(prs.slides):
        slide_content = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    text = para.text.strip()
                    if text:
                        slide_content.append(text)
        if slide_content:
            slides_text.append(f"[Diapositive {i+1}]\n" + "\n".join(slide_content))
    return "\n\n".join(slides_text)


def parse_eml(filepath: str) -> str:
    """
    Extrait le texte d'un fichier mail (.eml).
    Récupère l'expéditeur, le destinataire, le sujet, la date
    et le corps du message texte (ignore les pièces jointes).

    Paramètre :
    - filepath : chemin complet vers le fichier .eml

    Retourne :
    - contenu structuré du mail (str)
    """
    with open(filepath, "r", encoding="utf-8") as f:
        msg = email.message_from_file(f)

    headers = f"De : {msg.get('From', '')}\n"
    headers += f"À : {msg.get('To', '')}\n"
    headers += f"Sujet : {msg.get('Subject', '')}\n"
    headers += f"Date : {msg.get('Date', '')}\n"

    body = ""
    if msg.is_multipart():
        for part in msg.walk():
            if part.get_content_type() == "text/plain":
                body = part.get_payload(decode=True).decode("utf-8", errors="ignore")
                break
    else:
        body = msg.get_payload(decode=True).decode("utf-8", errors="ignore")

    return f"{headers}\n{body}"


def parse_file(filepath: str) -> dict:
    """
    Router de parsing — détecte le format selon l'extension
    et appelle le bon parser.

    Paramètre :
    - filepath : chemin complet vers le fichier à parser

    Retourne un dict :
    {
      "content"  : texte extrait (str),
      "source"   : nom du fichier (str),
      "format"   : extension du fichier (str),
      "error"    : message d'erreur si le parsing a échoué (str ou None)
    }
    """
    filename = os.path.basename(filepath)
    ext = os.path.splitext(filename)[1].lower()

    parsers = {
        ".txt":  parse_txt,
        ".docx": parse_docx,
        ".pptx": parse_pptx,
        ".eml":  parse_eml,
    }

    if ext not in parsers:
        return {
            "content": "",
            "source": filename,
            "format": ext,
            "error": f"Format non supporté : {ext}"
        }

    try:
        content = parsers[ext](filepath)
        return {
            "content": content,
            "source": filename,
            "format": ext,
            "error": None
        }
    except Exception as e:
        return {
            "content": "",
            "source": filename,
            "format": ext,
            "error": str(e)
        }


def parse_folder(folder_path: str) -> list[dict]:
    """
    Parse tous les fichiers supportés dans un dossier.
    Ignore les fichiers dont le format n'est pas supporté.

    Paramètre :
    - folder_path : chemin vers le dossier contenant les fichiers à parser

    Retourne :
    - liste de dicts produits par parse_file(), un par fichier parsé avec succès
    """
    SUPPORTED = {".txt", ".docx", ".pptx", ".eml"}
    results = []

    for filename in os.listdir(folder_path):
        ext = os.path.splitext(filename)[1].lower()
        if ext in SUPPORTED:
            filepath = os.path.join(folder_path, filename)
            result = parse_file(filepath)
            if result["error"]:
                print(f"Erreur parsing {filename} : {result['error']}")
            else:
                print(f"Parsé : {filename} ({len(result['content'])} caractères) [{ext}]")
                results.append(result)

    return results


# ── Génération de documentation ───────────────────────────────────────────────

# Prompt de structuration — instruit le LLM à produire un document
# propre à partir de fragments hétérogènes.
# Règles clés :
# - Signaler les lacunes avec [À COMPLÉTER] plutôt que d'inventer
# - Signaler les contradictions plutôt que de choisir arbitrairement
# - Produire un document Markdown structuré avec titres et sections
GENERATION_PROMPT = """Tu es un assistant de documentation pour une PME.
À partir des fragments de documents fournis, génère un document structuré en Markdown.

Règles strictes :
1. Utilise UNIQUEMENT les informations présentes dans les fragments fournis.
2. Si une information est manquante, indique [À COMPLÉTER] à l'endroit concerné.
3. Si deux sources se contredisent, indique [CONTRADICTION : source1 dit X, source2 dit Y].
4. Structure le document avec des titres Markdown (##, ###).
5. Ne fabrique aucune information absente des sources.
6. Indique les sources utilisées en fin de document."""


def generate_document(parsed_files: list[dict], doc_title: str) -> str:
    """
    Génère un document structuré en Markdown à partir de fichiers parsés.
    Utilise le LLM local (Ollama) ou cloud (API Mistral) selon LLM_MODE.

    Paramètres :
    - parsed_files : liste de dicts produits par parse_folder()
    - doc_title    : titre souhaité pour le document généré

    Retourne :
    - document structuré en Markdown (str)
    - message d'erreur si la génération a échoué (str)

    Note sur la validation :
    Ce document est un brouillon — il doit être relu et validé par
    un humain avant d'être indexé dans la base RAG.
    """
    # Assemblage du contexte : tous les fragments avec leur source
    context = "\n\n".join([
        f"[Source : {f['source']} — format {f['format']}]\n{f['content']}"
        for f in parsed_files
    ])

    prompt = f"""Titre du document à générer : {doc_title}

Fragments sources :
{context}

Génère le document structuré en Markdown."""

    if LLM_MODE == "cloud":
        return _generate_cloud(prompt)
    else:
        return _generate_local(prompt)


def _generate_local(prompt: str) -> str:
    """
    Génère via Ollama en local.
    Timeout élevé (3 min) car la génération d'un document complet
    est plus longue qu'une simple réponse RAG.
    """
    payload = {
        "model": OLLAMA_MODEL,
        "messages": [
            {"role": "system", "content": GENERATION_PROMPT},
            {"role": "user", "content": prompt}
        ],
        "stream": False
    }
    try:
        response = requests.post(
            f"{OLLAMA_BASE_URL}/api/chat",
            json=payload,
            timeout=180
        )
        if response.status_code == 200:
            return response.json()["message"]["content"]
        return f"Erreur Ollama : {response.status_code}"
    except requests.exceptions.ConnectionError:
        return "Erreur : Ollama n'est pas disponible."
    except requests.exceptions.Timeout:
        return "Erreur : délai dépassé. Essaie en mode cloud."
    
def _generate_cloud(prompt: str) -> str:
    """
    Génère via l'API Mistral cloud.
    Plus rapide que le mode local pour les documents longs.
    """
    if not MISTRAL_API_KEY:
        return "Erreur : clé API Mistral manquante dans le fichier .env"

    headers = {
        "Authorization": f"Bearer {MISTRAL_API_KEY}",
        "Content-Type": "application/json"
    }
    payload = {
        "model": MISTRAL_MODEL,
        "messages": [
            {"role": "system", "content": GENERATION_PROMPT},
            {"role": "user", "content": prompt}
        ]
    }
    try:
        response = requests.post(
            "https://api.mistral.ai/v1/chat/completions",
            headers=headers,
            json=payload,
            timeout=60
        )
        if response.status_code == 200:
            return response.json()["choices"][0]["message"]["content"]
        return f"Erreur API Mistral : {response.status_code}"
    except requests.exceptions.Timeout:
        return "Erreur : délai dépassé sur l'API Mistral."